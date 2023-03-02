import urllib.request
import sys
import os
from PIL import Image
import pandas as pd
from pptx import Presentation
from pptx.util import Inches


diretorio = 'C:/Users/tacilio.pereira/Downloads/bookdeFotos/2_fotos_book'

# Abre o arquivo CSV em modo de leitura
df = pd.read_excel('C:/Users/tacilio.pereira/Downloads/bookdeFotos/2_fotos.xlsb')

# Obtém a coluna desejada como um objeto Series
coluna1 = df.loc[:, 'URL_FOTO_1']
coluna2 = df.loc[:, 'URL_FOTO_2']
# Obtém o número de linhas no DataFrame
num_linhas = df.shape[0]
print('Iniciando o Download das Imagens')
print(f'Salvando no diretorio {diretorio}')
for i in range(num_linhas):
    for valor1, valor2 in zip(coluna1, coluna2):
        foto1 = f"foto{i}.jpg"
        foto2 = f"fotos{i}.jpg"
        url1 = f'{valor1}'
        url2 = f'{valor2}'
        try:
            caminho_arquivo1 = os.path.join(diretorio, foto1)
            caminho_arquivo2 = os.path.join(diretorio, foto2)
            img1 = urllib.request.urlretrieve(url1, caminho_arquivo1)
            img2 = urllib.request.urlretrieve(url2, caminho_arquivo2)
        except:
            erro = sys.exc_info()
            print("Ocorreu um erro:", erro)

print('FIM')

# Cria uma apresentação em branco
ppt = Presentation()

# Lista todos os arquivos no diretório
arquivos = os.listdir(diretorio)

# Filtra apenas os arquivos com extensão .jpg ou .png
imagens = [arquivo for arquivo in arquivos if arquivo.endswith('.jpg') or arquivo.endswith('.png')]

# Carrega cada imagem e processa uma por vez
for imagem in imagens:
    # Abre a imagem usando a biblioteca Pillow
    caminho_imagem = os.path.join(diretorio, imagem)
    print(caminho_imagem)
    img = Image.open(caminho_imagem)

    # Adiciona um slide em branco à apresentação
    slide = ppt.slides.add_slide(ppt.slide_layouts[6])
    # Insere uma imagem no slide
    left = Inches(1)
    top = Inches(2)
    slide.shapes.add_picture(caminho_imagem, left, top, height=Inches(4))

    # Fecha a imagem
    img.close()

