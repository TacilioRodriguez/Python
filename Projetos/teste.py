import urllib.request
import sys
import os
from PIL import Image
import pandas as pd
from pptx import Presentation
from pptx.util import Inches
import win32com.client

diretorio = 'C:/Users/tacilio.pereira/Downloads/bookdeFotos/2_fotos_book'

# Abre o arquivo CSV em modo de leitura
df = pd.read_excel('C:/Users/tacilio.pereira/Downloads/bookdeFotos/2_fotos.xlsb')

# Obtém a coluna desejada como um objeto Series
coluna1 = df.loc[:, 'URL_FOTO_1']
coluna2 = df.loc[:, 'URL_FOTO_2']
# Obtém o número de linhas no DataFrame
num_linhas = df.shape[0]
print('Iniciando o Download das Imagens')
print(f'Salvando no diretorio {diretorio}')

# Cria duas listas vazias para armazenar o nome das fotos
fotos1 = []
fotos2 = []

for i in range(num_linhas):
    valor1 = coluna1[i]
    valor2 = coluna2[i]
    url1 = f'{valor1}'
    url2 = f'{valor2}'
    try:
        foto1 = f"foto{i}.jpg"
        foto2 = f"fotos{i}.jpg"
        caminho_arquivo1 = os.path.join(diretorio, foto1)
        caminho_arquivo2 = os.path.join(diretorio, foto2)
        img1 = urllib.request.urlretrieve(url1, caminho_arquivo1)
        img2 = urllib.request.urlretrieve(url2, caminho_arquivo2)
        fotos1.append(foto1)
        fotos2.append(foto2)
    except:
        erro = sys.exc_info()
        print("Ocorreu um erro:", erro)



# Cria uma apresentação em branco
ppt = Presentation()

# Lista todos os arquivos no diretório
arquivos = os.listdir(diretorio)

# Filtra apenas os arquivos com extensão .jpg ou .png
imagens = [arquivo for arquivo in arquivos if arquivo.endswith('.jpg') or arquivo.endswith('.png')]

# Carrega cada imagem e processa uma por vez
for imagem in imagens:
    # Abre a imagem usando a biblioteca Pillow
    caminho_imagem = os.path.join(diretorio, imagem)
    print(caminho_imagem)
    img = Image.open(caminho_imagem)

    # Adiciona um slide em branco à apresentação
    slide = ppt.slides.add_slide(ppt.slide_layouts[6])
    # Insere uma imagem no slide
    left = Inches(1)
    top = Inches(2)
    slide.shapes.add_picture(caminho_imagem, left, top, height=Inches(4))

    # Fecha a imagem
    img.close()


# Cria uma instância do PowerPoint
powerpoint = win32com.client.Dispatch("Powerpoint.Application")

# Abre um arquivo de apresentação existente
presentation = powerpoint.Presentations.Open("C:/Users/tacilio.pereira/Downloads/bookdeFotos/2_fotos_book/Apresentação1.pptx")

# Seleciona o primeiro slide
slide = presentation.Slides(1)

print(os.listdir(diretorio))
print((os.path.join(diretorio, foto1)))
#Insere a primeira imagem no slide
left1 = Inches(1)
top1 = Inches(2)
slide.Shapes.AddPicture(os.path.join(diretorio, foto1), False, True, left1, top1, Inches(4), Inches(4))